# -*- coding: utf-8 -*-
"""
SC Fuels — AI-Driven Growth Blueprint
PowerPoint deck builder (python-pptx) with full Technijian branding + speaker notes.
Source content: Clients/SCF/_research.md  |  Brand: assets/brand-tokens.json
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = "c:/vscode/tech-branding/tech-branding/"
ASSETS = BASE + "assets/"
DIAG = BASE + "Clients/SCF/diagrams/"
OUT = BASE + "Clients/SCF/SC-Fuels-AI-Growth-Blueprint.pptx"

# ---------- Brand palette ----------
BLUE   = RGBColor(0x00, 0x6D, 0xB6)
ORANGE = RGBColor(0xF6, 0x7D, 0x4B)
TEAL   = RGBColor(0x1E, 0xAA, 0xC8)
CHART  = RGBColor(0xCB, 0xDB, 0x2D)
GREY   = RGBColor(0x59, 0x59, 0x5B)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
NEARBLK= RGBColor(0x2D, 0x2D, 0x2D)
OFFWHT = RGBColor(0xF8, 0xF9, 0xFA)
LTGREY = RGBColor(0xE9, 0xEC, 0xEF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x28, 0xA7, 0x45)
CRIT   = RGBColor(0xCC, 0x00, 0x00)
FONT   = "Open Sans"

# AUTHENTIC Technijian logo files (verified against technijian.com 2026-06-01).
# The logos/png/technijian-logo-* set are AI-regenerated approximations — DO NOT use.
LOGO_LIGHT = ASSETS + "Technijian Logo 2.png"            # authentic transparent full-color (light bg)
LOGO_DARK  = ASSETS + "Technijian Logo - white text.png" # authentic transparent white wordmark (dark bg)
LOGO_REV   = ASSETS + "Technijian Logo - white text.png" # authentic white wordmark (dark solid)

EMU_IN = 914400
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def slide():
    return prs.slides.add_slide(BLANK)

def _noline(shp):
    shp.line.fill.background()

def rect(s, x, y, w, h, color, line_color=None, line_w=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is None:
        _noline(sp)
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def rrect(s, x, y, w, h, color, line_color=None, line_w=1.0, radius=0.08):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is None:
        _noline(sp)
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, space_after=None, line_spacing=None):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, italic) run tuples,
       OR a single tuple for one-run paragraph."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space_after is not None: p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing is not None: p.line_spacing = line_spacing
        rl = para if isinstance(para, list) else [para]
        for (t, sz, col, bold, *rest) in rl:
            ital = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.italic = ital
            r.font.name = FONT
    return tb

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def pic_fit(s, path, box_x, box_y, box_w, box_h, align="center", valign="middle"):
    from PIL import Image
    iw, ih = Image.open(path).size
    r = min(box_w/iw, box_h/ih)
    w, h = iw*r, ih*r
    if align == "center": x = box_x + (box_w - w)/2
    elif align == "left": x = box_x
    else: x = box_x + (box_w - w)
    if valign == "middle": y = box_y + (box_h - h)/2
    elif valign == "top": y = box_y
    else: y = box_y + (box_h - h)
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return x, y, w, h

def header(s, eyebrow, title, accent=BLUE):
    rect(s, 0, 0, SW, 0.14, accent)                      # top accent strip
    rect(s, 0, 0.14, SW, 0.04, ORANGE)                   # thin orange underline
    # subtle brand logo top-right (authentic full-color, ratio 4.78)
    s.shapes.add_picture(LOGO_LIGHT, Inches(SW-2.0), Inches(0.46), width=Inches(1.6))
    txt(s, 0.7, 0.52, 9.5, 0.3, [[(eyebrow.upper(), 12, ORANGE, True)]])
    # title held to a SINGLE line (26pt, width clears the logo) so it never collides with the underline
    txt(s, 0.7, 0.84, 10.3, 0.62, [[(title, 26, BLUE, True)]])
    rect(s, 0.72, 1.62, 1.5, 0.045, ORANGE)              # title underline accent

def footer(s, n, dark=False):
    col = RGBColor(0xB8,0xB8,0xC4) if dark else RGBColor(0x9A,0x9A,0xA2)
    line = RGBColor(0x3A,0x3A,0x50) if dark else LTGREY
    rect(s, 0.7, 7.06, SW-1.4, 0.012, line)
    txt(s, 0.7, 7.12, 8.5, 0.3, [[("AI-Driven Growth Blueprint  ·  Prepared for SC Fuels  ·  Confidential", 8.5, col, False)]])
    txt(s, SW-4.3, 7.12, 3.0, 0.3, [[("technijian.com", 8.5, col, False)]], align=PP_ALIGN.RIGHT)
    txt(s, SW-1.0, 7.12, 0.3, 0.3, [[(str(n), 8.5, col, True)]], align=PP_ALIGN.RIGHT)

def bg(s, color):
    rect(s, -0.06, -0.06, SW+0.12, SH+0.12, color)

# ============================================================ SLIDE 1 — COVER
s = slide()
bg(s, DARK)
# left accent column of brand blocks
for i,c in enumerate([BLUE, TEAL, ORANGE, CHART]):
    rect(s, 0, i*(SH/4), 0.22, SH/4, c)
# decorative faint band
rect(s, 0, 5.55, SW, 0.02, RGBColor(0x33,0x33,0x4d))
s.shapes.add_picture(LOGO_DARK, Inches(0.95), Inches(0.95), width=Inches(4.0))
txt(s, 1.0, 2.75, 11.4, 0.4, [[("AI-DRIVEN GROWTH BLUEPRINT", 15, TEAL, True)]])
txt(s, 0.98, 3.2, 11.6, 1.5,
    [[("Defend share. Integrate acquisitions.", 40, WHITE, True)],
     [("Out-modernize the West's fuel market.", 40, WHITE, True)]],
    line_spacing=1.04)
txt(s, 1.0, 5.0, 11.0, 0.5,
    [[("A practical roadmap to grow SC Fuels with AI — built on Pilot's scale, not against your field team.", 15, RGBColor(0xC9,0xCb,0xD6), False, True)]])
txt(s, 1.0, 5.95, 11.4, 0.4,
    [[("Prepared for SC Fuels Leadership", 13, WHITE, True),
      ("      |      June 1, 2026      |      Confidential", 13, RGBColor(0x9A,0x9A,0xB0), False)]])
# (tagline is already part of the authentic logo mark above — no redundant line)
notes(s, "OPEN / FRAME (60–90 sec). Thank them for the time. Set the tone: this is a working blueprint, not a hard sell — a practical roadmap for where AI moves the needle for SC Fuels.\n\nThree things to land up front:\n1) We see SC Fuels as the MARKET LEADER — the West Coast commercial arm of Pilot Company (Berkshire-backed since Nov 2021), and the active consolidator (Downs Energy Oct 2025, Lutz Petroleum Dec 2025). We are NOT pitching you as a scrappy underdog.\n2) Every number in here is ESTIMATED and conservative until discovery — we'll be explicit about that.\n3) AI SUPPORTS your field reps and Pilot's supply scale — it does not replace the relationship-driven sales motion. That boundary is a trust signal; say it out loud.\n\nDO NOT cite a revenue figure (aggregator data is unreliable for a Pilot subsidiary). Use gallons/scale instead. Only CEO Steven Greinke is a verified name — keep other leadership names out.")

# ============================================================ SLIDE 2 — AGENDA
s = slide()
bg(s, WHITE)
header(s, "Today's conversation", "What we'll cover", BLUE)
items = [
 ("01", "Where SC Fuels stands today", "Market leader, Pilot-backed — and under-represented online", BLUE),
 ("02", "The market & regulatory landscape", "Consolidation, energy transition, LCFS / RIN / IFTA burden", TEAL),
 ("03", "Where the growth lives", "A hybrid engine: searchable demand + named accounts", ORANGE),
 ("04", "The competitive white space", "Nobody in the West has claimed the AI story — yet", BLUE),
 ("05", "How AI transforms the growth engine", "Get Found · Account Intelligence · Internal Automation", TEAL),
 ("06", "Investment, ROI & the roadmap", "Published rates, scoped in discovery · 90/180/365", ORANGE),
]
y = 1.95
for num, t, sub, c in items:
    rrect(s, 0.7, y, 0.62, 0.62, c, radius=0.18)
    txt(s, 0.7, y, 0.62, 0.62, [[(num, 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.5, y+0.02, 10.8, 0.4, [[(t, 16.5, NEARBLK, True)]])
    txt(s, 1.5, y+0.36, 10.8, 0.3, [[(sub, 12, GREY, False)]])
    y += 0.81
footer(s, 2)
notes(s, "AGENDA (30 sec). Quick roadmap of the next ~30–40 minutes. This is meant to be interactive — please interrupt with questions and corrections; half the value today is you telling us where our estimates are wrong.\n\nWe'll end on QUICK WINS you can start this week with zero commitment, and the handful of numbers we need from you to firm up the ROI.")

# ============================================================ SLIDE 3 — WHO SC FUELS IS TODAY
s = slide()
bg(s, WHITE)
header(s, "Where SC Fuels stands today", "A market leader, backed by Pilot scale", BLUE)
kpis = [
 ("1930", "Founded", "Signal Oil distributorship", BLUE),
 ("~900", "Employees", "Across 15+ Western states", TEAL),
 ("11,000+", "Commercial customers", "Family business → Fortune 500", ORANGE),
 ("1B+", "Gallons / year", "Fuel + lubricants + DEF", BLUE),
 ("15+", "Western states", "SoCal HQ + regional branches", TEAL),
 ("Pilot", "Berkshire-backed", "Acquired Nov 2021 · the consolidator", ORANGE),
]
cw, ch, gx, gy = 3.78, 1.95, 0.28, 0.32
x0, y0 = 0.7, 2.0
for i,(num,lab,sub,c) in enumerate(kpis):
    cx = x0 + (i%3)*(cw+gx); cy = y0 + (i//3)*(ch+gy)
    card = rrect(s, cx, cy, cw, ch, WHITE, line_color=LTGREY, line_w=1.0, radius=0.06)
    rect(s, cx, cy, 0.10, ch, c)                       # left accent
    txt(s, cx+0.32, cy+0.22, cw-0.5, 0.7, [[(num, 40, c, True)]])
    txt(s, cx+0.34, cy+1.05, cw-0.5, 0.35, [[(lab, 14.5, NEARBLK, True)]])
    txt(s, cx+0.34, cy+1.42, cw-0.5, 0.4, [[(sub, 11, GREY, False)]])
footer(s, 3)
notes(s, "THE REFRAME — most important strategic slide (2–3 min).\n\nThe headline: SC Fuels is the West's market-leading fuel & lubricant distributor, and since November 2021 it's been the West Coast commercial arm of PILOT COMPANY, majority-owned by BERKSHIRE HATHAWAY. It's also the consolidator — it acquired Downs Energy (Oct 2025) and Lutz Petroleum (Dec 2025). So the balance sheet, supply infrastructure, and CFN/cardlock network behind SC Fuels are top-tier.\n\nWhy this matters for the plan: the digital presence today reads like a small regional independent — that mismatch IS the opportunity. We're not trying to make you bigger; we're making the digital + AI presence finally match the operation.\n\nGUARDRAILS: Do NOT cite revenue (aggregators conflict $394M vs $650M — unreliable for a Pilot subsidiary). The 'third-generation family-owned' copy still on the website is stale/legacy — we'll flag it later as a liability to fix. The ~1.4B gallon figure is a pre-2021 number; say 'more than a billion gallons.' Numbers are scale anchors, not guarantees.")

# ============================================================ SLIDE 4 — MARKET POSITION
s = slide()
bg(s, WHITE)
header(s, "The strategic frame", "Three jobs AI does for a market leader", BLUE)
cards = [
 ("DEFEND", "Defend share", "Out-service the regionals and lock in the 11,000-account base with predictive retention and a better self-serve experience.", BLUE),
 ("INTEGRATE", "Integrate acquisitions", "Turn the Downs & Lutz customer books into clean onboarding and cross-sell — account intelligence instead of post-deal churn.", TEAL),
 ("OUT-MODERNIZE", "Out-modernize disruptors", "Adopt Booster's UX and AI bar — app-grade portal, answer-engine visibility — without the venture burn, on Pilot's supply.", ORANGE),
]
cw, gx = 3.86, 0.28
x0, y0, ch = 0.7, 2.15, 3.7
for i,(tag,t,body,c) in enumerate(cards):
    cx = x0 + i*(cw+gx)
    rrect(s, cx, y0, cw, ch, OFFWHT, line_color=LTGREY, radius=0.05)
    rect(s, cx, y0, cw, 0.14, c)
    rrect(s, cx+0.3, y0+0.42, 1.9, 0.46, c, radius=0.2)
    txt(s, cx+0.3, y0+0.42, 1.9, 0.46, [[(tag, 11.5, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx+0.3, y0+1.15, cw-0.6, 0.5, [[(t, 19, NEARBLK, True)]])
    txt(s, cx+0.3, y0+1.75, cw-0.6, 1.7, [[(body, 13, GREY, False)]], line_spacing=1.25)
txt(s, 0.7, 6.15, SW-1.4, 0.5,
    [[("The boundary: ", 13, ORANGE, True),
      ("AI supports your field reps and Pilot's supply scale — it doesn't replace the relationship that wins fleet accounts.", 13, NEARBLK, False, True)]])
footer(s, 4)
notes(s, "STRATEGIC FRAME (2 min). For a market leader, AI does three jobs — defend, integrate, out-modernize. Walk each card briefly.\n\nDEFEND: ~14% of fleet fuel spend is lost to fraud/theft industry-wide; declining-gallon accounts are an early churn signal a model can flag for rep outreach.\nINTEGRATE: you're actively buying companies — every acquired customer book is either clean cross-sell or messy churn. AI account-intelligence + a knowledge graph makes it the former.\nOUT-MODERNIZE: Booster is the digital-native benchmark — app-first mobile fueling with a named 'Kasparov' AI routing engine. The play is to adopt the good UX without the startup cash burn, backed by Pilot's supply.\n\nClose on the BOUNDARY line at the bottom — this is the single most important trust signal in the whole deck. We are not promising to replace your sales force; we make them faster and better-armed.")

# ============================================================ SLIDE 5 — BUSINESS MODEL (diagram)
s = slide()
bg(s, WHITE)
header(s, "Where you make money", "The distribution business model", BLUE)
pic_fit(s, DIAG+"model.png", 0.7, 1.85, SW-1.4, 5.0, valign="top")
footer(s, 5)
notes(s, "BUSINESS MODEL (1.5 min). Quick orientation so everyone shares the same map. Supply (refiners, Chevron-branded, renewable) flows IN to SC Fuels, out through FOUR channels — bulk/bobtail delivery, mobile on-site (wet) fueling, cardlock/fleet cards, and unbranded wholesale + lubricants/DEF/renewables — to SIX verticals: construction, ag, manufacturing, food manufacturing, municipal/government, waste.\n\nKey point: this is a thin-margin, distributor-in-the-middle business — so growth and retention both matter, and small efficiency gains compound. The 'AI acceleration' band shows where we add leverage.\n\nCARDLOCK NUANCE: SC Fuels doesn't own a national cardlock brand — it runs proprietary sites and rides the CFN and Pacific Pride (Corpay) networks, same as rivals. So 'we have cardlock' differentiates nobody — the differentiation is service, data, portal UX, and supply. Don't oversell cardlock as a moat.")

# ============================================================ SLIDE 6 — INDUSTRY & REGULATORY
s = slide()
bg(s, WHITE)
header(s, "The landscape", "A consolidating, regulated, transitioning market", BLUE)
# left column: market forces
txt(s, 0.7, 1.9, 5.7, 0.35, [[("MARKET FORCES", 13, BLUE, True)]])
forces = [
 ("Heavy consolidation", "30+ distributor deals in 2024 — and you're on the buy side."),
 ("Margin pressure + energy transition", "Efficiency and depot EV charging encroach on medium-duty fuel demand."),
 ("Renewable diesel (R99) is the offset", "A drop-in substitute; LCFS credits drive the economics. You already sell it."),
 ("Telematics convergence", "Cards + telematics + tank sensors merging; ~14% of fleet fuel spend lost to fraud."),
]
y = 2.3
for t, b in forces:
    rect(s, 0.72, y+0.06, 0.12, 0.12, ORANGE)
    txt(s, 0.98, y, 5.4, 0.32, [[(t, 13.5, NEARBLK, True)]])
    txt(s, 0.98, y+0.30, 5.4, 0.5, [[(b, 11.5, GREY, False)]], line_spacing=1.15)
    y += 1.02
# right column: regulatory table
txt(s, 6.7, 1.9, 6.0, 0.35, [[("COMPLIANCE BURDEN  —  AI-REDUCTION TARGETS", 13, BLUE, True)]])
regs = [
 ("CARB / LCFS", "CI target 22.75% (Jul 2025); credit accounting", "HIGH", ORANGE),
 ("RFS / RINs", "RIN tracking, RVO, attestation", "HIGH", ORANGE),
 ("IFTA", "Multi-state fuel tax — ~88 hrs/yr/fleet", "HIGH", ORANGE),
 ("UST rules", "Single-wall tank closure by Dec 2025", "MED", TEAL),
 ("DOT / Hazmat", "Driver/Hazmat certs, BOL/manifests", "MED", TEAL),
]
ry = 2.35
rect(s, 6.7, ry, 6.0, 0.02, LTGREY)
for name, req, sev, sc in regs:
    rect(s, 6.7, ry+0.08, 6.0, 0.78, OFFWHT)
    txt(s, 6.85, ry+0.14, 2.2, 0.3, [[(name, 12.5, NEARBLK, True)]])
    txt(s, 6.85, ry+0.45, 4.9, 0.35, [[(req, 10.5, GREY, False)]], line_spacing=1.05)
    rrect(s, 11.55, ry+0.24, 1.0, 0.36, sc, radius=0.25)
    txt(s, 11.55, ry+0.24, 1.0, 0.36, [[(sev, 10, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.92
footer(s, 6)
notes(s, "LANDSCAPE (2 min). Two halves: the market forces (left) and the regulatory burden (right).\n\nMarket forces — consolidation (you're the acquirer, good position), margin/energy-transition pressure (efficiency + medium-duty EV nibble at demand), renewable diesel as the margin OFFSET (you already sell R99 — room to lead), and telematics/fraud convergence.\n\nThe right side is the AI thesis's fastest proof: LCFS + RIN + IFTA are heavy, recurring, error-prone, and CARB/LCFS is UNIQUELY WESTERN — so your 15-state footprint makes it especially acute. IFTA alone is ~88 hours/year per fleet of manual work. These are bounded, painful, and clear — the perfect Tier-3 'prove value fast' wedge.\n\nDISCOVERY ASK: we need your actual hours and cost spent on LCFS, RIN, and IFTA today — that's what anchors the compliance-automation ROI. Flag that you'll come back to it on the questions slide.")

# ============================================================ SLIDE 7 — WHERE GROWTH LIVES
s = slide()
bg(s, WHITE)
header(s, "Where the growth lives", "A hybrid engine — not shotgun marketing", BLUE)
pools = [
 ("Searchable demand", "\"bulk diesel delivery [city]\", \"cardlock near me\", \"DEF delivery\", \"renewable diesel supplier\" — real demand you under-harvest today.", TEAL),
 ("Named accounts", "Large fleets, construction & ag operations — account intelligence + per-account dossiers that arm your field reps across 15 states.", BLUE),
 ("Municipal RFPs", "Bid-response automation turns days into minutes — countering the missing MBE/WBE diversity wedge with speed and Pilot scale.", ORANGE),
 ("Renewable / ESG", "R99 opportunity intelligence — match LCFS economics to the accounts and regions where renewable diesel wins.", GREEN),
]
cw, ch, gx, gy = 5.85, 1.85, 0.28, 0.28
x0, y0 = 0.7, 1.95
for i,(t,b,c) in enumerate(pools):
    cx = x0 + (i%2)*(cw+gx); cy = y0 + (i//2)*(ch+gy)
    rrect(s, cx, cy, cw, ch, OFFWHT, line_color=LTGREY, radius=0.05)
    rect(s, cx, cy, cw, 0.13, c)
    txt(s, cx+0.32, cy+0.33, cw-0.6, 0.4, [[(t, 17, c, True)]])
    txt(s, cx+0.32, cy+0.85, cw-0.64, 0.9, [[(b, 12.5, GREY, False)]], line_spacing=1.22)
txt(s, 0.7, 6.5, SW-1.4, 0.4,
    [[("Honest framing: ", 12.5, ORANGE, True),
      ("AI supports the rep relationship and Pilot's supply — it doesn't replace them. We harvest demand you already have; we don't manufacture noise.", 12.5, NEARBLK, False, True)]])
footer(s, 7)
notes(s, "WHERE GROWTH LIVES (2 min) — the honesty slide.\n\nThis is a HYBRID engine leaning account-based — exactly right for a B2B field-sales business. Four pools: (1) searchable demand you under-harvest, (2) named high-value accounts for your reps, (3) municipal RFPs where speed wins, (4) renewable/ESG where margin lives.\n\nDo NOT say 'fill your funnel' or 'shotgun marketing' — that would be the wrong pitch and they'd know it. The honesty IS the differentiator.\n\nTHE ONE STRUCTURAL DISADVANTAGE to name proactively: SC Fuels lacks the MBE/WBE diversity-spend status that Pinnacle (woman-owned) and Ramos (Hispanic-owned) use to win municipal set-aside contracts. We do NOT pretend to fix that — we counter it with bid-response SPEED (RFP automation), service depth, and Pilot scale. Naming the weakness honestly builds trust.")

# ============================================================ SLIDE 8 — CUSTOMER (diagram)
s = slide()
bg(s, WHITE)
header(s, "Who buys", "The SC Fuels customer — volume × margin", BLUE)
pic_fit(s, DIAG+"personas.png", 0.7, 1.8, SW-1.4, 5.05, valign="top")
footer(s, 8)
notes(s, "THE CUSTOMER (1.5 min). Four PRIMARY buyer personas — Fleet Operations Director (trucking/logistics), Construction Equipment/Ops Manager, Agricultural Ops Manager, and Municipal Fleet & Procurement Officer — plus two EMERGING ones: the Sustainability/ESG fleet buyer and the Lubricants & DEF industrial buyer.\n\nPlotted on Account Fuel Volume (x) vs Margin / Strategic Value (y). Notice the emerging personas (ESG renewable, lubricants/DEF) sit higher on margin — they deserve disproportionate attention even at lower volume. Municipal is low-margin but HIGH strategic value (sticky, multi-year, but bid-driven and diversity-weighted).\n\nThese are research-derived candidates — ask them to validate or correct in discovery. If they tell us their best accounts look different, we re-plot.")

# ============================================================ SLIDE 9 — COMPETITIVE (diagram)
s = slide()
bg(s, WHITE)
header(s, "The white space", "Competitive positioning — scale × AI maturity", BLUE)
x,y,w,h = pic_fit(s, DIAG+"competitive.png", 0.7, 1.8, 8.1, 5.0, align="left", valign="top")
# callout panel on the right
px = 9.1
rrect(s, px, 2.0, 3.55, 4.5, DARK, radius=0.05)
rect(s, px, 2.0, 3.55, 0.13, ORANGE)
txt(s, px+0.3, 2.35, 3.0, 0.4, [[("THE OPEN CORNER", 13, ORANGE, True)]])
wsp = [
 "Mansfield & World Kinect lead on content + data.",
 "Booster leads on UX + named AI.",
 "You already beat the dated regionals (Hunt & Sons, Ramos).",
 "Nobody in your Western backyard has claimed the AI / answer-engine / compliance story.",
]
yy = 2.85
for w_ in wsp:
    rect(s, px+0.3, yy+0.07, 0.1, 0.1, TEAL)
    txt(s, px+0.52, yy, 2.85, 0.8, [[(w_, 11.5, WHITE, False)]], line_spacing=1.18)
    yy += 0.86
txt(s, px+0.3, 6.05, 3.0, 0.35, [[("First mover wins compounding inbound.", 11.5, CHART, True, True)]])
footer(s, 9)
notes(s, "COMPETITIVE WHITE SPACE (2 min) — the 'why now' slide.\n\nEight named competitors. The digital leaders are World Kinect/Flyers and Mansfield (content + fuel-data portals); Booster is the digital-native (app + named 'Kasparov' AI). The dated regionals — Hunt & Sons (Phillips 66-owned, fragmented domains), Ramos — you already out-class digitally.\n\nThe white space (right panel): the SCALE + AI corner is uncontested in the West. Nobody in your backyard has claimed the AI / answer-engine / compliance-automation narrative. First mover compounds — AI answer engines and Google reward the authority that publishes first.\n\nProof point to mention: Star Oilco (tiny Portland distributor) has best-in-class content/SEO for its size — proof that content works even in this 'boring' category. If a small shop can rank, a Pilot-backed leader can dominate.")

# ============================================================ SLIDE 10 — DIGITAL AUDIT
s = slide()
bg(s, WHITE)
header(s, "The honest audit", "A market leader, under-represented online", BLUE)
finds = [
 ("AI-blank & content-light", "Thin vs Mansfield / World Kinect (data) and Booster (named AI + app).", CRIT),
 ("No location / branch pages", "~15 branches, zero per-city landing pages — real local demand left on the table.", ORANGE),
 ("Stale \"family-owned\" copy", "Contradicts the Pilot/Berkshire reality — a credibility + SEO liability, and a selling point unused.", CRIT),
 ("Light social footprint", "~5K LinkedIn followers for a ~900-person, billion-gallon leader.", TEAL),
 ("Solid self-serve portal", "Rebates, e-receipts, transaction controls — a strong base to EXTEND, not rebuild.", GREEN),
]
y = 2.0
for t, b, c in finds:
    rrect(s, 0.7, y, SW-1.4, 0.86, OFFWHT, line_color=LTGREY, radius=0.06)
    rect(s, 0.7, y, 0.12, 0.86, c)
    txt(s, 1.05, y+0.13, 4.6, 0.6, [[(t, 15, NEARBLK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 5.8, y+0.13, 6.9, 0.62, [[(b, 12, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y += 0.98
footer(s, 10)
notes(s, "DIGITAL AUDIT (1.5–2 min). The on-ramp: the quality and scale of the operation is FAR ahead of how it presents online and in AI answers. This isn't criticism — it's unclaimed upside.\n\nHit the two red items hardest:\n- 'Family-owned' copy is ACTIVELY costing you: it's a credibility/SEO liability AND it leaves the Pilot-scale selling point on the table. Fixing it is a same-day quick win.\n- No location pages = the single biggest missed-capture opportunity for a 15-branch operation. Real people are googling 'bulk diesel delivery [city]' and not finding you.\n\nEnd positive (the green item): you already have a real self-serve portal. We EXTEND it toward the Booster UX bar — we don't rip and replace. That lowers cost and risk.")

# ============================================================ SLIDE 11 — CAPABILITY PROOF
s = slide()
bg(s, WHITE)
header(s, "Technijian capability", "Five proven builds — mapped to SC Fuels", BLUE)
builds = [
 ("Multi-Agent SEO + AEO", "Own \"bulk fuel delivery [city]\", renewable, cardlock & DEF queries in Google AND AI answers — plus per-branch location pages.", BLUE),
 ("AI Lead Gen — Account-Based", "Named commercial / fleet / construction / municipal target lists + per-account dossiers for reps across 15 states.", TEAL),
 ("AI Document Intelligence", "Municipal & large-fleet fuel RFP response — days to minutes. The speed answer to the diversity-wedge gap.", ORANGE),
 ("AI-Native Custom App (My Dev)", "Self-serve portal + AI chatbot + LCFS / RIN / IFTA automation + predictive tank monitoring.", BLUE),
 ("Local SEO + Reputation", "Google Business Profiles + reviews across ~15 branches — lifts local capture and AI-citation trust.", TEAL),
]
cw, gx = 3.86, 0.28
# 3 on top, 2 on bottom centered
positions = [(0.7,2.0),(4.84,2.0),(8.98,2.0),(2.77,4.45),(6.91,4.45)]
ch = 2.25
for (cx,cy),(t,b,c) in zip(positions, builds):
    rrect(s, cx, cy, cw, ch, WHITE, line_color=LTGREY, radius=0.06)
    rect(s, cx, cy, cw, 0.13, c)
    txt(s, cx+0.28, cy+0.36, cw-0.55, 0.7, [[(t, 14.5, NEARBLK, True)]], line_spacing=1.0)
    txt(s, cx+0.28, cy+1.05, cw-0.55, 1.1, [[(b, 11.5, GREY, False)]], line_spacing=1.2)
footer(s, 11)
notes(s, "CAPABILITY PROOF (2 min). These are real, productized Technijian builds — mapped to your situation.\n\nLANGUAGE DISCIPLINE (important): present these as PROVEN CAPABILITIES and productized services we're mapping to SC Fuels — not as 'we already built this exact thing for a fuel distributor.' We have NOT built an AI engine for a fuel distributor before; say so plainly if asked. Honesty here protects credibility on every other number.\n\nThe one to emphasize for THEIR pain: AI Document Intelligence for RFPs — it's the speed answer to the diversity-spend disadvantage. And the custom app is PHASED — discovery first, not a big-bang build.\n\nIf Derek's original question was about SEO specifically, anchor on builds #1 and #5 — that's the warm entry point.")

# ============================================================ SLIDE 12 — AI ENGINE (diagram)
s = slide()
bg(s, WHITE)
header(s, "The engine", "How AI transforms the growth engine", BLUE)
pic_fit(s, DIAG+"architecture.png", 0.7, 1.8, SW-1.4, 5.0, valign="top")
footer(s, 12)
notes(s, "THE AI GROWTH ENGINE (2 min). Three columns: GET FOUND (SEO / AEO / local + location pages), ACCOUNT INTELLIGENCE & OUTBOUND (lead gen + RFP automation + rep dossiers), and INTERNAL AUTOMATION (compliance + predictive tank monitoring + demand/route forecasting).\n\nSEQUENCING is the message: lead with Tier 1 top-line GROWTH (exploits the AI-blank gap, gets you found and arms the reps), prove value FAST with Tier 3 COMPLIANCE (bounded, painful, measurable — LCFS/IFTA), then Tier 2 retains and expands (forecasting, tank monitoring, churn).\n\nThe boundary stays on screen the whole time: every box SUPPORTS a person — the rep, the compliance manager, the dispatcher. AI is the force-multiplier, not the replacement.")

# ============================================================ SLIDE 13 — INVESTMENT & ROI
s = slide()
bg(s, WHITE)
header(s, "Investment", "Built on published rates — scoped in discovery", BLUE)
# table: real published list rates where they exist; TBD (set in discovery) otherwise
rows = [
 ("My SEO", "Local + AEO + location/branch pages", "$500–$1,500/mo", True),
 ("My AI Lead Gen", "Account intelligence & outbound", "$1,499–$6,999/mo", True),
 ("My AI — Fractional AI Advisor", "Program leadership & governance", "TBD", False),
 ("My Dev — Custom build", "Portal · chatbot · LCFS/IFTA automation", "TBD", False),
 ("My Dev — Managed app services", "Hosting · monitoring · support", "TBD", False),
 ("My AI — Executive Workshop", "Leadership alignment + roadmap", "TBD", False),
]
tx, tw = 0.7, 7.45
ty = 2.0
# header row
rect(s, tx, ty, tw, 0.42, BLUE)
txt(s, tx+0.2, ty, 4.3, 0.42, [[("SERVICE", 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx+4.6, ty, 2.65, 0.42, [[("INVESTMENT", 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
ty += 0.42
for i,(svc,cov,val,real) in enumerate(rows):
    bgc = WHITE if i%2==0 else OFFWHT
    rect(s, tx, ty, tw, 0.6, bgc)
    txt(s, tx+0.2, ty+0.10, 5.0, 0.3, [[(svc, 11.5, NEARBLK, True)]])
    txt(s, tx+0.2, ty+0.34, 5.0, 0.24, [[(cov, 9.5, GREY, False)]])
    if real:
        txt(s, tx+4.6, ty, 2.65, 0.6, [[(val, 12.5, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    else:
        txt(s, tx+4.6, ty, 2.65, 0.6, [[(val, 12, GREY, True, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    ty += 0.6
# total row
rect(s, tx, ty, tw, 0.52, DARK)
txt(s, tx+0.2, ty, 4.0, 0.52, [[("YEAR-1 TOTAL", 12.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx+3.5, ty, 3.75, 0.52, [[("TBD — finalized after discovery", 12, CHART, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
ty += 0.62
txt(s, tx, ty, tw, 0.5,
    [[("Published Technijian list rates. SEO/Lead-Gen tier + custom-build scope set in discovery; Lead Gen adds one-time setup ($2,500 / $5,000 / $15,000).", 8.5, GREY, False, True)]],
    line_spacing=1.1)
# ROI panel — method, not invented multiples
rx = 8.45
rrect(s, rx, 2.0, 4.15, 4.74, DARK, radius=0.05)
rect(s, rx, 2.0, 4.15, 0.14, ORANGE)
txt(s, rx+0.32, 2.3, 3.55, 0.4, [[("HOW WE'LL MEASURE ROI", 12.5, ORANGE, True)]])
txt(s, rx+0.32, 2.74, 3.55, 0.6,
    [[("Modeled on gross profit — a thin-margin fuel business, never on revenue.", 10.5, RGBColor(0xC9,0xCb,0xD6), False, True)]], line_spacing=1.15)
method = [
 "New commercial accounts  ×  gross profit per account",
 "+  Retained / expanded volume on the 11,000-account base",
 "−  Program investment",
]
yy = 3.72
for mline in method:
    rect(s, rx+0.32, yy+0.04, 0.1, 0.1, TEAL)
    txt(s, rx+0.54, yy, 3.3, 0.55, [[(mline, 10.5, WHITE, False)]], line_spacing=1.1)
    yy += 0.62
txt(s, rx+0.32, 5.66, 3.55, 0.6,
    [[("Targets are set from YOUR baseline — GP/account, new-accounts/mo, churn — captured in discovery.", 10, RGBColor(0xC9,0xCb,0xD6), False)]], line_spacing=1.12)
txt(s, rx+0.32, 6.4, 3.55, 0.3, [[("Illustrative until then — no number we can't back.", 10.5, CHART, True, True)]])
footer(s, 13)
notes(s, "INVESTMENT (3 min) — truthful framing. Lead with the honesty: today we show the PUBLISHED Technijian list rates where they apply, and we do NOT put a number on anything we'd have to guess.\n\nWhat's real on this slide: My SEO published tiers run $500–$1,500/mo (Tier 1 Website/Hosting up to Tier 5 Topical Viral Videos, each tier includes everything below it). My AI Lead Gen is $1,499 Starter / $3,499 Professional / $6,999 Enterprise per month, plus one-time setup ($2,500 / $5,000 / $15,000). Which SEO tier and which Lead Gen tier SC Fuels needs is a SCOPE decision we make together in discovery — so it's a range here, not a pinned number.\n\nWhat's TBD: the Fractional AI Advisor rate, the custom build (portal + chatbot + LCFS/IFTA), managed app services, and the workshop — we do not have a defensible number for these until we scope SC Fuels' actual systems, gallons, account count, and compliance hours. The Year-1 total is therefore TBD by design. Say that plainly — it signals we quote what we can stand behind.\n\nROI: we do NOT promise a multiple. We show the METHOD — gross profit per new account + retained/expanded volume, minus investment — and we set the targets from THEIR baseline (GP/account, new-accounts/mo, churn) in discovery. If asked 'what's the return,' answer: 'That's exactly what discovery sizes — we won't hand you a number we can't back.'\n\nLAND-AND-EXPAND still applies: start small (a published-rate SEO tier + a location-page pilot + the workshop), prove value, then scope the build. The ask today is the first step, not a total.")

# ============================================================ SLIDE 14 — ROADMAP (diagram)
s = slide()
bg(s, WHITE)
header(s, "The plan", "Implementation roadmap — 90 / 180 / 365", BLUE)
pic_fit(s, DIAG+"timeline.png", 0.7, 1.95, SW-1.4, 4.5, valign="middle")
footer(s, 14)
notes(s, "ROADMAP (1.5 min). Enterprise, regulated, multi-stakeholder, multi-state — so the pace is deliberate, not rushed.\n\n90 DAYS — Foundation: Executive AI Workshop, SEO/AEO + the first location pages, Google Business Profiles, ONE compliance wedge (LCFS or IFTA) to prove value fast, and the first account-target lists.\n180 DAYS — Scale: RFP automation live, portal/chatbot, outbound at volume, churn/at-risk model on the cardlock base.\n365 DAYS — Optimize: full compliance suite, predictive tank monitoring, demand/route forecasting, and an acquisition-integration playbook for the next Downs/Lutz.\n\nThe message: you see value inside 90 days (the compliance wedge + location pages), but the big bets are sequenced so nothing is bet-the-farm.")

# ============================================================ SLIDE 15 — QUICK WINS
s = slide()
bg(s, WHITE)
header(s, "Start this week", "Five quick wins — zero commitment", ORANGE)
qw = [
 ("Build one location page", "For your largest branch — a template that proves the local-capture gap."),
 ("Claim Google Business Profiles", "Across all ~15 branches — fast, visible local-search lift."),
 ("Publish one LCFS / renewable FAQ", "Own the AI answer competitors haven't claimed yet."),
 ("Refresh the ownership story", "Align the website with the Pilot-backed reality — liability into selling point."),
 ("Capture 5 account testimonials", "Proof assets for account-based selling."),
]
y = 2.05
for i,(t,b) in enumerate(qw):
    rrect(s, 0.7, y, SW-1.4, 0.84, OFFWHT, line_color=LTGREY, radius=0.07)
    rrect(s, 0.92, y+0.19, 0.46, 0.46, ORANGE, radius=0.22)
    txt(s, 0.92, y+0.19, 0.46, 0.46, [[(str(i+1), 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.6, y+0.13, 5.0, 0.6, [[(t, 15, NEARBLK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 6.7, y+0.13, 6.0, 0.6, [[(b, 12, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    y += 0.96
footer(s, 15)
notes(s, "QUICK WINS (1.5 min) — the most important slide for momentum.\n\nThese are deliberately ZERO-COMMITMENT, low-risk proofs. They let us demonstrate value BEFORE any meaningful spend — which de-risks the whole engagement for them.\n\nThe ask: pick ONE or TWO today and we'll just do them. The fastest visible combo is #1 + #2 (a location page + Google Business Profiles) — it produces real local-search movement in weeks. #4 (ownership-copy refresh) is almost free and removes an active liability.\n\nDon't leave this slide without a concrete 'which one should we start with?' — that converts a presentation into an engagement.")

# ============================================================ SLIDE 16 — DISCOVERY QUESTIONS
s = slide()
bg(s, WHITE)
header(s, "Calibrate the plan", "What we need from you to firm up the numbers", BLUE)
qs = [
 "Current gallons/yr, active accounts, new-accounts/month, and sales headcount by territory",
 "Churn rate on cardlock and commercial accounts (baselines the retention story)",
 "Hours and cost spent on LCFS + RIN + IFTA today (the compliance-ROI anchors)",
 "CRM and billing/back-office systems (the integration surface for AI)",
 "Branch count + exact states, and renewable-diesel as a % of total volume",
 "Who owns marketing/web today — and the relationship to Pilot's central marketing",
 "The engagement sponsor — SC Fuels leadership, or routed through Pilot?",
]
y = 1.95
for i,q in enumerate(qs):
    c = ORANGE if i==len(qs)-1 else BLUE
    rect(s, 0.78, y+0.07, 0.14, 0.14, c)
    bold = i==len(qs)-1
    txt(s, 1.12, y, 11.4, 0.5, [[(q, 13.5, NEARBLK if not bold else c, bold)]], line_spacing=1.1)
    y += 0.66
txt(s, 0.7, 6.7, SW-1.4, 0.4,
    [[("Every projection in this deck is an estimate until we answer these together. ", 12, GREY, False, True),
      ("The sponsor question changes the whole motion.", 12, ORANGE, True)]])
footer(s, 16)
notes(s, "DISCOVERY (1.5 min) — signals rigor and respect. Frame it as: 'every number you saw is an estimate until we fill these in together.'\n\nIf there's time, capture answers LIVE — even partial answers sharpen the next version of the plan.\n\nThe LAST question (highlighted) matters most: is the buyer SC Fuels' own marketing/sales leadership, or is this routed through Pilot's central marketing? That single answer changes the entire sales motion, the budget authority, and who else needs to be in the room. If you only get one answer today, get that one.\n\nSecondary priorities: LCFS/RIN/IFTA hours (anchors compliance ROI) and GP-per-account (anchors growth ROI).")

# ============================================================ SLIDE 17 — NEXT / CTA
s = slide()
bg(s, DARK)
for i,c in enumerate([ORANGE, CHART, TEAL, BLUE]):
    rect(s, (SW/4)*i, SH-0.18, SW/4, 0.18, c)
s.shapes.add_picture(LOGO_DARK, Inches(0.95), Inches(0.8), width=Inches(3.4))
txt(s, 0.98, 2.0, 11.0, 0.4, [[("WHERE WE GO FROM HERE", 14, TEAL, True)]])
txt(s, 0.95, 2.45, 11.4, 0.8, [[("Three small steps — no big commitment", 34, WHITE, True)]])
steps = [
 ("1", "Pick 1–2 quick wins", "We start this week — zero commitment."),
 ("2", "60-minute discovery", "Firm up the numbers behind the plan."),
 ("3", "Executive AI Workshop", "Align leadership on the roadmap."),
]
x = 0.98
for num,t,b in steps:
    rrect(s, x, 3.65, 3.7, 1.55, RGBColor(0x24,0x24,0x3e), radius=0.06)
    rrect(s, x+0.28, 3.9, 0.5, 0.5, ORANGE, radius=0.22)
    txt(s, x+0.28, 3.9, 0.5, 0.5, [[(num, 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x+0.95, 3.95, 2.6, 0.5, [[(t, 15, WHITE, True)]], line_spacing=1.0)
    txt(s, x+0.3, 4.55, 3.2, 0.6, [[(b, 11.5, RGBColor(0xC2,0xC4,0xD2), False)]], line_spacing=1.15)
    x += 3.92
# contact bar
rrect(s, 0.98, 5.55, 11.37, 0.92, RGBColor(0x12,0x12,0x22), radius=0.05)
txt(s, 1.3, 5.7, 7.5, 0.3, [[("Ravi Jain", 15, WHITE, True), ("   Founder & CEO, Technijian", 12, RGBColor(0xB9,0xBd,0xCc), False)]])
txt(s, 1.3, 6.06, 10.5, 0.32, [[("rjain@technijian.com    |    949.379.8499    |    technijian.com    |    ", 12, TEAL, False),
                                ("technology as a solution", 12, ORANGE, True, True)]])
notes(s, "CLOSE (1 min) — on momentum, not pressure.\n\nThe ask is small and concrete: pick a quick win and book the 60-minute discovery. The Executive AI Workshop is the natural first paid step (priced once we scope it in discovery) — low-risk, high-alignment, and it gets leadership in one room to shape the roadmap.\n\nReiterate the core message one last time: we're here to help the MARKET LEADER extend its lead. AI supports your people and Pilot's scale — it doesn't replace them. We move at your pace, every number gets validated in discovery, and you can start this week with zero commitment.\n\nThen stop talking and ask: 'Which quick win do you want us to start with?' Let them answer.\n\nLogistics reminder: this engagement is being run for Bryan Burkhart; Derek Bettencourt opened the door asking about SEO. If Derek is in the room, anchor the warm entry on the SEO/AEO + location-page story (slides 11 #1/#5 and quick wins #1/#2).")

prs.save(OUT)
print("SAVED:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
